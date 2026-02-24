"""
Generate a professional PowerPoint presentation for CodeMix NLP project.
Creates an animated, visually appealing PPTX with diagrams, tables, and charts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)        # Deep navy background
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)         # Card background
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)     # Bright cyan-blue
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)    # Green
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)      # Red
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)   # Orange/amber
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)   # Purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)      # Muted text
MEDIUM_GRAY = RGBColor(0x64, 0x74, 0x8B)      # Subtle text
GRADIENT_START = RGBColor(0x06, 0xB6, 0xD4)   # Teal
GRADIENT_END = RGBColor(0x81, 0x67, 0xF3)     # Violet

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_entrance_animation(slide, shape, delay_ms=0, duration_ms=500):
    """Add a fade-in entrance animation to a shape."""
    timing = slide.element.find(qn('p:timing'))
    if timing is None:
        timing = slide.element.makeelement(qn('p:timing'), {})
        slide.element.append(timing)

    tn_lst = timing.find(qn('p:tnLst'))
    if tn_lst is None:
        tn_lst = timing.makeelement(qn('p:tnLst'), {})
        timing.append(tn_lst)

    par = tn_lst.find(qn('p:par'))
    if par is None:
        par = tn_lst.makeelement(qn('p:par'), {})
        tn_lst.append(par)
        c_tn_root = par.makeelement(qn('p:cTn'), {
            'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
        })
        par.append(c_tn_root)
        child_tn_lst = c_tn_root.makeelement(qn('p:childTnLst'), {})
        c_tn_root.append(child_tn_lst)
        seq = child_tn_lst.makeelement(qn('p:seq'), {
            'concurrent': '1', 'nextAc': 'seek'
        })
        child_tn_lst.append(seq)
        c_tn_seq = seq.makeelement(qn('p:cTn'), {
            'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'
        })
        seq.append(c_tn_seq)
        seq_child = c_tn_seq.makeelement(qn('p:childTnLst'), {})
        c_tn_seq.append(seq_child)
    else:
        c_tn_root = par.find(qn('p:cTn'))
        child_tn_lst = c_tn_root.find(qn('p:childTnLst'))
        seq = child_tn_lst.find(qn('p:seq'))
        c_tn_seq = seq.find(qn('p:cTn'))
        seq_child = c_tn_seq.find(qn('p:childTnLst'))

    # Create animation par node
    anim_par = seq_child.makeelement(qn('p:par'), {})
    seq_child.append(anim_par)
    anim_ctn = anim_par.makeelement(qn('p:cTn'), {
        'id': str(len(seq_child) + 2), 'fill': 'hold'
    })
    anim_par.append(anim_ctn)
    stCondLst = anim_ctn.makeelement(qn('p:stCondLst'), {})
    anim_ctn.append(stCondLst)
    cond = stCondLst.makeelement(qn('p:cond'), {'delay': '0'})
    stCondLst.append(cond)

    inner_child = anim_ctn.makeelement(qn('p:childTnLst'), {})
    anim_ctn.append(inner_child)

    inner_par = inner_child.makeelement(qn('p:par'), {})
    inner_child.append(inner_par)
    inner_ctn = inner_par.makeelement(qn('p:cTn'), {
        'id': str(len(seq_child) + 3), 'presetID': '10', 'presetClass': 'entr',
        'presetSubtype': '0', 'fill': 'hold', 'nodeType': 'withEffect'
    })
    inner_par.append(inner_ctn)

    st2 = inner_ctn.makeelement(qn('p:stCondLst'), {})
    inner_ctn.append(st2)
    cond2 = st2.makeelement(qn('p:cond'), {'delay': str(delay_ms)})
    st2.append(cond2)

    inner_child2 = inner_ctn.makeelement(qn('p:childTnLst'), {})
    inner_ctn.append(inner_child2)

    # AnimEffect (fade)
    anim_effect = inner_child2.makeelement(qn('p:animEffect'), {
        'transition': 'in', 'filter': 'fade'
    })
    inner_child2.append(anim_effect)
    ae_ctn = anim_effect.makeelement(qn('p:cBhvr'), {})
    anim_effect.append(ae_ctn)
    ae_ctn2 = ae_ctn.makeelement(qn('p:cTn'), {
        'id': str(len(seq_child) + 4), 'dur': str(duration_ms)
    })
    ae_ctn.append(ae_ctn2)
    tgt_el = ae_ctn.makeelement(qn('p:tgtEl'), {})
    ae_ctn.append(tgt_el)
    sp_tgt = tgt_el.makeelement(qn('p:spTgt'), {'spid': str(shape.shape_id)})
    tgt_el.append(sp_tgt)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # Remove existing bullets
        for bu in pPr.findall(qn('a:buChar')):
            pPr.remove(bu)
        for bu in pPr.findall(qn('a:buNone')):
            pPr.remove(bu)
        pPr.append(buChar)
        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color[0]:02X}{bullet_color[1]:02X}{bullet_color[2]:02X}'})
        buClr.append(srgbClr)
        for bc in pPr.findall(qn('a:buClr')):
            pPr.remove(bc)
        pPr.append(buClr)

    return txBox


def add_table(slide, left, top, rows, cols, col_widths, data, header_color=ACCENT_BLUE):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          sum(col_widths), Inches(0.4) * rows)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Style header row
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = data[0][j]
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0x6E)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = ACCENT_BLUE
            paragraph.font.bold = True
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.CENTER

    # Style data rows
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j]
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x15, 0x1F, 0x34) if i % 2 == 1 else CARD_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = "Segoe UI"
                paragraph.alignment = PP_ALIGN.CENTER

    # Remove table borders
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is None:
                    tcPr = tc.makeelement(qn('a:tcPr'), {})
                    tc.append(tcPr)
                ln = tcPr.makeelement(qn(f'a:{border_name}'), {'w': '6350'})
                for old in tcPr.findall(qn(f'a:{border_name}')):
                    tcPr.remove(old)
                solidFill = ln.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '1E293B'})
                solidFill.append(srgb)
                ln.append(solidFill)
                tcPr.append(ln)

    return table_shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total):
    """Add slide number to bottom-right."""
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)


# Total slides count
TOTAL_SLIDES = 20


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Decorative top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

# University name
add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.5),
             "VISHWAKARMA UNIVERSITY, PUNE", font_size=16, color=LIGHT_GRAY, bold=False)

# Thin line
add_accent_line(slide, Inches(1), Inches(1.4), Inches(2.5), ACCENT_BLUE)

# Main title
t1 = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "CodeMix NLP", font_size=52, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
             "Sarcasm & Misinformation Detection in Hinglish Text", font_size=28, color=ACCENT_BLUE, bold=False)

# Description
add_text_box(slide, Inches(1), Inches(3.9), Inches(8), Inches(0.7),
             "A Multi-Task NLP System using XLM-RoBERTa with SHAP Explainability",
             font_size=16, color=LIGHT_GRAY)

# Info boxes at bottom
info_items = [
    ("Subject", "Natural Language Processing"),
    ("Program", "B.Tech — 3rd Year, Sem 2"),
    ("Student", "Sumeet Sangwan"),
    ("Year", "2025")
]
for i, (label, value) in enumerate(info_items):
    x = Inches(1) + Inches(2.8) * i
    card = add_shape(slide, x, Inches(5.2), Inches(2.5), Inches(1.2), CARD_BG)
    add_text_box(slide, x + Inches(0.2), Inches(5.35), Inches(2.1), Inches(0.4),
                 label, font_size=11, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(5.7), Inches(2.1), Inches(0.5),
                 value, font_size=14, color=WHITE, bold=False)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Agenda", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

toc_left = [
    "01  Introduction & Problem Statement",
    "02  Objectives",
    "03  NLP Concepts & Techniques",
    "04  System Architecture",
    "05  Technology Stack",
    "06  Dataset Details",
    "07  XLM-RoBERTa Model",
    "08  Multi-Task Learning",
    "09  NLP Pipeline Walkthrough",
    "10  Sarcasm Detection Engine",
]
toc_right = [
    "11  Misinformation Detection Engine",
    "12  Training Pipeline",
    "13  SHAP Explainability",
    "14  Backend Implementation",
    "15  Frontend Implementation",
    "16  API Documentation",
    "17  Live Demo & Results",
    "18  Key Achievements",
    "19  Future Enhancements",
    "20  Thank You & Q&A",
]

for i, item in enumerate(toc_left):
    y = Inches(1.6) + Inches(0.48) * i
    num_part = item[:2]
    text_part = item[4:]
    add_text_box(slide, Inches(1), y, Inches(0.5), Inches(0.4),
                 num_part, font_size=15, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.5), y, Inches(5), Inches(0.4),
                 text_part, font_size=15, color=LIGHT_GRAY)

for i, item in enumerate(toc_right):
    y = Inches(1.6) + Inches(0.48) * i
    num_part = item[:2]
    text_part = item[4:]
    add_text_box(slide, Inches(7), y, Inches(0.5), Inches(0.4),
                 num_part, font_size=15, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(7.5), y, Inches(5), Inches(0.4),
                 text_part, font_size=15, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Introduction & Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Introduction & Problem Statement", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# Left: The Problem
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
             "The Challenge", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4), [
    "750M+ Indian internet users communicate in code-mixed languages",
    "Hinglish = Hindi + English mixed in one sentence",
    "Standard English NLP models FAIL on code-mixed text",
    "Sarcasm relies on cultural context, tone & emojis",
    "WhatsApp misinformation spreads in Hinglish/Hindi",
    "No existing tool handles both tasks simultaneously",
], font_size=14, bullet_color=ACCENT_RED)

# Right: Example
add_shape(slide, Inches(6.8), Inches(1.4), Inches(5.7), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.4),
             "Sarcasm Example", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5.1), Inches(0.7),
             '"Haan bilkul, Modi ji ne toh desh ka\nbohot vikas kar diya 🙄"',
             font_size=16, color=ACCENT_ORANGE, bold=False, font_name="Consolas")
add_text_box(slide, Inches(7.1), Inches(3.0), Inches(5.1), Inches(0.5),
             "Words say POSITIVE → Intent is NEGATIVE", font_size=13, color=LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.7), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(4.4), Inches(5), Inches(0.4),
             "Misinformation Example", font_size=18, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(7.1), Inches(5.0), Inches(5.1), Inches(0.7),
             '"5G towers se corona failta hai,\nscientists ne prove kiya — share karo!"',
             font_size=16, color=ACCENT_ORANGE, bold=False, font_name="Consolas")
add_text_box(slide, Inches(7.1), Inches(5.8), Inches(5.1), Inches(0.5),
             "Fake claim + False authority + Urgency", font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Objectives  
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Project Objectives", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

objectives = [
    ("🎯", "Multilingual NLP Model", "Build a model that understands code-mixed\nHinglish, Tanglish, and English text", ACCENT_BLUE),
    ("🔀", "Dual-Task Detection", "Detect BOTH sarcasm AND misinformation\nsimultaneously with a single model", ACCENT_GREEN),
    ("📊", "Training Dataset", "Generate 13,792 synthetic training samples\nacross 3 languages and 8 categories", ACCENT_ORANGE),
    ("🔍", "Explainability (SHAP)", "Show token-level contributions — WHY the\nmodel made each prediction", ACCENT_PURPLE),
    ("🌐", "Full-Stack Web App", "Real-time analysis, dashboards, history\ntracking with modern UI", ACCENT_BLUE),
    ("✅", "High Accuracy", "Achieve 100% accuracy on diverse test set\nincluding political sarcasm & health misinfo", ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4) * col
    y = Inches(1.5) + Inches(2.8) * row

    card = add_shape(slide, x, y, Inches(3.7), Inches(2.4), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=28, color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.3), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.3), Inches(0.9),
                 desc, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: NLP Concepts & Techniques
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "NLP Concepts & Techniques Used", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

nlp_data = [
    ["Technique", "Implementation", "Purpose"],
    ["Transformer Architecture", "Self-Attention + Multi-Head", "Contextual word representations"],
    ["XLM-RoBERTa-large", "355M params, 100 languages", "Multilingual text understanding"],
    ["Transfer Learning", "Freeze first 6/24 layers", "Reuse pre-trained knowledge"],
    ["Multi-Task Learning", "Shared encoder + 2 heads", "Joint sarcasm + misinfo detection"],
    ["SentencePiece Tokenization", "Subword BPE tokens", "Handle multiple scripts"],
    ["Emoji Normalization", "28 emojis → semantic tokens", "Preserve emoji meaning"],
    ["Language Detection", "Unicode + word frequency", "Detect hinglish/tanglish/english"],
    ["Occlusion SHAP", "Mask each token, measure Δ", "Token-level explainability"],
    ["Attention Visualization", "Last-layer attention weights", "Model interpretability"],
    ["Rule-Based NLP", "13 sarcasm + 8 misinfo rules", "Linguistic pattern matching"],
    ["Data Augmentation", "Template + prefix/suffix + synonym", "Dataset diversity"],
]

add_table(slide, Inches(0.5), Inches(1.4), len(nlp_data), 3,
          [Inches(3.2), Inches(4.2), Inches(4.9)], nlp_data)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: System Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "System Architecture", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

# User layer
add_shape(slide, Inches(4.5), Inches(1.3), Inches(4.3), Inches(0.7), RGBColor(0x15, 0x50, 0x3C))
add_text_box(slide, Inches(4.5), Inches(1.35), Inches(4.3), Inches(0.6),
             "👤  User (Web Browser — localhost:3000)", font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow
add_shape(slide, Inches(6.5), Inches(2.05), Inches(0.3), Inches(0.4), ACCENT_BLUE)

# Frontend layer
add_shape(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.4), RGBColor(0x1E, 0x3A, 0x5F))
add_text_box(slide, Inches(2.8), Inches(2.55), Inches(7.7), Inches(0.4),
             "Frontend — Next.js 14 + React 18 + TypeScript", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
frontend_items = ["Homepage (Demo)", "Analyze Page", "Dashboard", "History", "TanStack Query + Axios"]
for i, item in enumerate(frontend_items):
    x = Inches(2.7) + Inches(1.6) * i
    add_text_box(slide, x, Inches(3.1), Inches(1.5), Inches(0.5),
                 item, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, Inches(5.8), Inches(3.9), Inches(1.5), Inches(0.4),
             "REST API (JSON)", font_size=11, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(6.5), Inches(4.2), Inches(0.3), Inches(0.3), ACCENT_ORANGE)

# Backend layer
add_shape(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.7), RGBColor(0x2D, 0x1B, 0x3D))
add_text_box(slide, Inches(0.9), Inches(4.55), Inches(11.4), Inches(0.4),
             "Backend — FastAPI (Python 3.11)", font_size=16, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Backend sub-boxes
be_boxes = [
    ("Middleware", "CORS → GZip →\nRate Limit → Logging", Inches(0.9), Inches(5.15), Inches(2.5)),
    ("API Router /api/v1", "/analyze | /analytics\n/models | /health", Inches(3.6), Inches(5.15), Inches(2.5)),
    ("ML Service (584 lines)", "Preprocess → Detect Lang\n→ Score Sarcasm (13 rules)\n→ Score Misinfo (8 rules)", Inches(6.3), Inches(5.15), Inches(3.1)),
    ("DB + Explain", "SQLite storage\nSHAP explainability", Inches(9.6), Inches(5.15), Inches(2.7)),
]

for title, desc, x, y, w in be_boxes:
    add_shape(slide, x, y, w, Inches(1.8), CARD_BG)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.35),
                 title, font_size=11, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(1.1),
                 desc, font_size=10, color=LIGHT_GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Technology Stack
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Technology Stack", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

# Backend stack
add_shape(slide, Inches(0.5), Inches(1.4), Inches(4), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(3.4), Inches(0.4),
             "⚙️  Backend (Python)", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(3.4), Inches(4.5), [
    "FastAPI 0.109.2 — Web framework",
    "PyTorch ≥2.0 — Deep learning",
    "Transformers 4.38 — HuggingFace",
    "SHAP 0.44 — Explainability",
    "Captum 0.7 — Interpretability",
    "SQLAlchemy 2.0 — ORM (async)",
    "Scikit-learn — ML metrics",
    "Structlog — JSON logging",
    "Pydantic 2.6 — Validation",
], font_size=12, bullet_color=ACCENT_BLUE)

# Frontend stack
add_shape(slide, Inches(4.8), Inches(1.4), Inches(4), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(5.1), Inches(1.55), Inches(3.4), Inches(0.4),
             "🎨  Frontend (TypeScript)", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(5.1), Inches(2.1), Inches(3.4), Inches(4.5), [
    "Next.js 14.1 — React framework",
    "React 18 — UI components",
    "TypeScript 5.3 — Type safety",
    "TanStack Query 5 — Data fetching",
    "Framer Motion 11 — Animations",
    "Tailwind CSS 3.4 — Styling",
    "Recharts 2.12 — Chart library",
    "Radix UI — Accessible primitives",
    "Zod 3.22 — Schema validation",
], font_size=12, bullet_color=ACCENT_GREEN)

# Infrastructure stack
add_shape(slide, Inches(9.1), Inches(1.4), Inches(3.7), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(9.4), Inches(1.55), Inches(3.1), Inches(0.4),
             "🏗️  Infrastructure", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(9.4), Inches(2.1), Inches(3.1), Inches(4.5), [
    "Docker Compose — 5 services",
    "PostgreSQL 15 — Production DB",
    "SQLite — Development DB",
    "Redis 7 — Response caching",
    "Celery 5.3 — Task queue",
    "Alembic — DB migrations",
    "Uvicorn — ASGI server",
    "GZip & Rate Limiting",
], font_size=12, bullet_color=ACCENT_ORANGE)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Dataset Details
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Dataset — 13,792 Synthetic Samples", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# Stats cards
stats = [
    ("13,792", "Total Samples", ACCENT_BLUE),
    ("3", "Languages", ACCENT_GREEN),
    ("8", "Categories", ACCENT_ORANGE),
    ("~76.3%", "Baseline F1", ACCENT_PURPLE),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.5) + Inches(3.15) * i
    add_shape(slide, x, Inches(1.4), Inches(2.9), Inches(1.2), CARD_BG)
    add_text_box(slide, x, Inches(1.5), Inches(2.9), Inches(0.6),
                 num, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.0), Inches(2.9), Inches(0.4),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Distribution table
dist_data = [
    ["Category", "Hinglish", "Tanglish", "English"],
    ["Sarcastic", "~3,750", "~625", "~1,000"],
    ["Non-Sarcastic", "~3,750", "~625", "~750"],
    ["Misinformation", "~3,125", "~480", "~750"],
    ["Credible", "~1,875", "—", "~750"],
]
add_table(slide, Inches(0.5), Inches(3.0), len(dist_data), 4,
          [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)], dist_data)

# Augmentation
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.4),
             "Data Augmentation Techniques", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0), [
    "Template Filling: Templates with {topic}, {person}, {claim} placeholders filled from 10+ variable pools",
    "Prefix Injection: Random prefixes (Bhai, Yaar, Da, Honestly, LOL)  |  Suffix Injection: Emoji/punctuation variations",
    "Word-Level Variation: Synonym substitution (bahut ↔ kaafi ↔ bohot, bhai ↔ yaar ↔ dost)",
], font_size=12, bullet_color=ACCENT_GREEN)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: XLM-RoBERTa Model
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "XLM-RoBERTa-large — Base Model", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# Key specs
specs = [
    ("355M", "Parameters", ACCENT_BLUE),
    ("100", "Languages", ACCENT_GREEN),
    ("1024", "Hidden Dim", ACCENT_ORANGE),
    ("24", "Transformer\nLayers", ACCENT_PURPLE),
]
for i, (num, label, color) in enumerate(specs):
    x = Inches(0.5) + Inches(3.15) * i
    add_shape(slide, x, Inches(1.4), Inches(2.9), Inches(1.3), CARD_BG)
    add_text_box(slide, x, Inches(1.5), Inches(2.9), Inches(0.6),
                 num, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.1), Inches(2.9), Inches(0.5),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Why XLM-RoBERTa
add_shape(slide, Inches(0.5), Inches(3.0), Inches(6), Inches(4), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(3.15), Inches(5.4), Inches(0.4),
             "Why XLM-RoBERTa for This Project?", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(3.7), Inches(5.4), Inches(3), [
    "Pre-trained on 2.5 TB CommonCrawl data",
    "Native support for Hindi, Tamil, English",
    "SentencePiece handles romanized Hindi (achha, bohot)",
    "Cross-lingual: similar words → nearby vectors",
    "MLM pre-training captures bidirectional context",
    "State-of-the-art on multilingual benchmarks",
], font_size=14, bullet_color=ACCENT_BLUE)

# Transfer Learning
add_shape(slide, Inches(6.8), Inches(3.0), Inches(6), Inches(4), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(3.15), Inches(5.4), Inches(0.4),
             "Transfer Learning Strategy", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(3.7), Inches(5.4), Inches(3), [
    "FROZEN: Embedding layer + Layers 1-6",
    "TRAINABLE: Layers 7-24 + Classification heads",
    "Lower layers = universal language features",
    "Upper layers = task-specific semantics",
    "Prevents catastrophic forgetting",
    "Reduces compute by ~40%",
], font_size=14, bullet_color=ACCENT_GREEN)

add_slide_number(slide, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Multi-Task Learning Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Multi-Task Learning Architecture", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# Input
add_shape(slide, Inches(3.5), Inches(1.4), Inches(6.3), Inches(0.7), RGBColor(0x1A, 0x35, 0x5C))
add_text_box(slide, Inches(3.5), Inches(1.45), Inches(6.3), Inches(0.6),
             "Input: Tokenized Text  [<s>, ▁haan, ▁bil, kul, ▁modi, ▁ji, ..., </s>]",
             font_size=13, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, Inches(6.3), Inches(2.1), Inches(0.7), Inches(0.3), "▼", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Shared Encoder
add_shape(slide, Inches(2.5), Inches(2.4), Inches(8.3), Inches(2.2), RGBColor(0x1E, 0x40, 0x6E))
add_text_box(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.4),
             "Shared XLM-RoBERTa Encoder", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), Inches(3.0), Inches(3.7), Inches(0.4),
             "🔒 Frozen: Embeddings + Layers 1-6", font_size=13, color=ACCENT_RED)
add_text_box(slide, Inches(6.8), Inches(3.0), Inches(3.7), Inches(0.4),
             "🔓 Trainable: Layers 7-24", font_size=13, color=ACCENT_GREEN)
add_text_box(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.5),
             "Output: [CLS] Token Representation (1024-dim vector)", font_size=14, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(4.1), Inches(8.3), Inches(0.4),
             "24 transformer layers • 16 attention heads • GELU activation", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Split arrows
add_text_box(slide, Inches(4.0), Inches(4.6), Inches(1.5), Inches(0.3), "▼", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.8), Inches(4.6), Inches(1.5), Inches(0.3), "▼", font_size=20, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# Sarcasm Head
add_shape(slide, Inches(1.5), Inches(4.9), Inches(5), Inches(2.2), RGBColor(0x15, 0x3A, 0x55))
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(5), Inches(0.4),
             "Sarcasm Classification Head", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.8), Inches(5.5), Inches(4.4), Inches(1.4),
             "Linear(1024 → 512) → GELU → LayerNorm\n→ Dropout(0.1) → Linear(512 → 2)\n→ Softmax → P(sarcastic)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Misinfo Head
add_shape(slide, Inches(6.8), Inches(4.9), Inches(5), Inches(2.2), RGBColor(0x3D, 0x15, 0x20))
add_text_box(slide, Inches(6.8), Inches(5.0), Inches(5), Inches(0.4),
             "Misinfo Classification Head", font_size=16, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.1), Inches(5.5), Inches(4.4), Inches(1.4),
             "Linear(1024 → 512) → GELU → LayerNorm\n→ Dropout(0.1) → Linear(512 → 2)\n→ Softmax → P(misinfo)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: NLP Pipeline Walkthrough
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "NLP Pipeline — Step by Step", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

pipeline_steps = [
    ("1", "Emoji\nNormalization", "🙄 → [EYEROLL]\n😂 → [LAUGH]\n28 emoji mappings", ACCENT_BLUE),
    ("2", "Text\nCleaning", "Remove URLs\nRemove @mentions\nStrip #hashtags\nLowercase", ACCENT_GREEN),
    ("3", "Language\nDetection", "Unicode ranges +\nRomanized word\nfrequency analysis", ACCENT_ORANGE),
    ("4", "Sarcasm\nScoring", "13 weighted rules\nPattern matching\nLexicon density\nContrast detection", ACCENT_PURPLE),
    ("5", "Misinfo\nScoring", "8 weighted rules\nUrgency language\nHealth claims\nFalse authority", ACCENT_RED),
    ("6", "Response\nAssembly", "Scores + Labels\nConfidence bands\nLanguage tag\nProcessing time", ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(pipeline_steps):
    x = Inches(0.3) + Inches(2.1) * i
    
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), Inches(1.5), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Connector line
    if i < 5:
        add_shape(slide, x + Inches(1.3), Inches(1.73), Inches(0.9), Inches(0.04), color)
    
    # Content card
    add_shape(slide, x, Inches(2.3), Inches(1.9), Inches(4.5), CARD_BG)
    add_text_box(slide, x + Inches(0.1), Inches(2.5), Inches(1.7), Inches(0.7),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.3), Inches(1.7), Inches(3.2),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Sarcasm Detection Engine
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Sarcasm Detection Engine — 13 Rules", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

sarcasm_data = [
    ["#", "Rule", "Weight", "Example Pattern"],
    ["1", "Structural Patterns (18 regex)", "0.35", "\"oh sure\", \"waah kya\", \"haan bilkul\""],
    ["2", "Sarcasm Word Density (~60 words)", "0.08/w", "waah, shabash, bilkul, genius, bravo"],
    ["3", "Sarcasm Emoji Detection", "0.20", "🙄 😏 😒 😂 👏 🤡"],
    ["4", "Strong Sarcasm Phrases", "0.25", "\"yeah right\", \"oh fantastic\", \"slow clap\""],
    ["5", "Contrast Detection", "0.30", "positive word + negative word in same text"],
    ["6", "Punctuation Irony", "0.15", "Multiple !!! or ??? or ... "],
    ["7", "Political Sarcasm", "0.30", "Political ref + development/progress claim"],
    ["8", "Negation Irony", "0.20", "\"not like it will work\", \"nahi karega\""],
    ["9", "Number + Positive Irony", "0.20", "\"50th meeting will definitely solve\""],
    ["10", "Exaggerated Praise", "0.25", "\"bohot vikas\", \"bahut achha kaam\""],
    ["11", "\"ne toh\" Pattern", "0.30", "\"Modi ji ne toh desh badal diya\""],
    ["12", "Mutual Exclusivity", "-0.5×", "High misinfo dampens sarcasm score"],
    ["13", "False Positive Dampening", "-20%", "Reduces score if only 1 weak signal"],
]

add_table(slide, Inches(0.3), Inches(1.3), len(sarcasm_data), 4,
          [Inches(0.5), Inches(3.8), Inches(1.2), Inches(6.8)], sarcasm_data)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Misinformation Detection Engine  
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Misinformation Detection Engine — 8 Rules", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3.5), ACCENT_RED)

misinfo_data = [
    ["#", "Rule", "Weight", "Example Pattern"],
    ["1", "Urgency Phrases (26 regex)", "0.30", "\"share before deleted\", \"forward to all\""],
    ["2", "Misinfo Word Density (~55 words)", "0.08/w", "cure, miracle, conspiracy, exposed, banned"],
    ["3", "Health Claims", "0.35", "\"cures cancer\", \"immunity boost\", \"home remedy\""],
    ["4", "False Authority", "0.25", "\"scientists proved\", \"doctors say\", \"NASA confirmed\""],
    ["5", "Conspiracy Language", "0.30", "\"hidden truth\", \"they don't want you to know\""],
    ["6", "ALL CAPS Ratio", "0.20", ">30% uppercase = sensationalist writing"],
    ["7", "Mutual Exclusivity", "-0.6×", "High sarcasm dampens misinformation score"],
    ["8", "Credibility Dampening", "-25%", "Reduces score if only 1 weak signal present"],
]

add_table(slide, Inches(0.5), Inches(1.3), len(misinfo_data), 4,
          [Inches(0.5), Inches(3.8), Inches(1.2), Inches(7)], misinfo_data)

# Hindi patterns box
add_shape(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(0.4),
             "Hindi/Hinglish Misinformation Patterns", font_size=18, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.5), [
    "\"share karo sabko\" — Urgency in Hindi",
    "\"scientists ne prove kiya\" — False authority",
    "\"ye sach hai forward karo\" — Claims truth",
    "\"5G se corona failta hai\" — Conspiracy",
], font_size=13, bullet_color=ACCENT_RED)
add_bullet_list(slide, Inches(6.8), Inches(5.4), Inches(5.5), Inches(1.5), [
    "\"cow urine cures COVID\" — Health myth",
    "\"delete hone se pehle\" — Create urgency",
    "\"government chupa rahi hai\" — Conspiracy",
    "\"100% guaranteed\" — False promise",
], font_size=13, bullet_color=ACCENT_ORANGE)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Training Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Training Pipeline", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

training_data = [
    ["Parameter", "Value", "Rationale"],
    ["Base Model", "xlm-roberta-large", "Best multilingual performance"],
    ["Learning Rate", "2e-5", "Standard for transformer fine-tuning"],
    ["Epochs", "5", "With early stopping (patience=2)"],
    ["Batch Size", "16 (effective 32)", "Gradient accumulation steps = 2"],
    ["Warmup Steps", "500", "Gradual LR increase"],
    ["Weight Decay", "0.01", "L2 regularization"],
    ["FP16", "Auto on CUDA", "2× faster training on GPU"],
    ["Frozen Layers", "Embed + Layers 1-6", "Preserve universal features"],
    ["Eval Steps", "Every 200 steps", "Monitor overall F1"],
    ["Save Strategy", "Best overall F1", "Max 3 checkpoints"],
]

add_table(slide, Inches(0.5), Inches(1.3), len(training_data), 3,
          [Inches(3), Inches(4), Inches(5.3)], training_data)

# Loss formula
add_shape(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
             "Multi-Task Loss Function", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
             "L = 0.5 × CrossEntropy(sarcasm) + 0.5 × CrossEntropy(misinfo)     |     Optimizer: AdamW     |     Scheduler: Linear warmup + decay",
             font_size=14, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER, font_name="Consolas")

add_slide_number(slide, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: SHAP Explainability
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "SHAP Explainability — Why Did the Model Decide?", font_size=30, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(4), ACCENT_BLUE)

# Algorithm
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(3.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.4),
             "Occlusion-Based SHAP Algorithm", font_size=17, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(2.6),
             "1. Get base prediction P(full_text)\n"
             "2. For each token tᵢ:\n"
             "      • Mask tᵢ → create text without it\n"
             "      • Get prediction P(text \\ tᵢ)\n"
             "      • Attribution = P(full) - P(masked)\n"
             "3. Positive score → pushes toward label\n"
             "4. Negative score → pushes against label",
             font_size=13, color=LIGHT_GRAY, font_name="Consolas")

# Example
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(3.5), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.4),
             "Example: Token Contributions", font_size=17, color=ACCENT_GREEN, bold=True)

tokens_example = [
    ("\"Haan\"", "+0.12", "Sarcastic opener", ACCENT_ORANGE),
    ("\"bilkul\"", "+0.15", "Strong sarcasm marker", ACCENT_RED),
    ("\"Modi ji\"", "+0.08", "Political reference", ACCENT_ORANGE),
    ("\"bohot\"", "+0.11", "Emphasis/exaggeration", ACCENT_ORANGE),
    ("\"vikas\"", "+0.14", "Positive in sarcastic ctx", ACCENT_RED),
    ("[EYEROLL]", "+0.22", "Strongest signal ★", ACCENT_RED),
]

for i, (token, score, desc, color) in enumerate(tokens_example):
    y = Inches(2.1) + Inches(0.4) * i
    add_text_box(slide, Inches(7.1), y, Inches(1.3), Inches(0.35),
                 token, font_size=11, color=WHITE, font_name="Consolas")
    add_text_box(slide, Inches(8.5), y, Inches(0.8), Inches(0.35),
                 score, font_size=11, color=color, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(9.5), y, Inches(3), Inches(0.35),
                 desc, font_size=11, color=LIGHT_GRAY)

# Visualization desc
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.4),
             "Frontend Visualization", font_size=17, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.85), Inches(5.5), Inches(1.2), [
    "SHAP Token Highlighting: Red (positive) / Green (negative)",
    "Opacity proportional to contribution magnitude",
    "Three-tabbed view: Sarcasm SHAP | Misinfo SHAP | Attention",
], font_size=13, bullet_color=ACCENT_PURPLE)
add_bullet_list(slide, Inches(6.8), Inches(5.85), Inches(5.5), Inches(1.2), [
    "Attention Heatmap: Canvas-rendered matrix (20×20)",
    "Self-attention from last transformer layer",
    "Shows which tokens attend to each other",
], font_size=13, bullet_color=ACCENT_BLUE)

add_slide_number(slide, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Backend Implementation
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Backend Implementation", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

# File structure
add_shape(slide, Inches(0.5), Inches(1.4), Inches(4), Inches(5.7), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(3.4), Inches(0.4),
             "📁 Project Structure", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.05), Inches(3.4), Inches(4.8),
             "backend/\n"
             "├── app/\n"
             "│   ├── main.py\n"
             "│   ├── api/v1/\n"
             "│   │   ├── analyze.py (4 routes)\n"
             "│   │   ├── analytics.py (3 routes)\n"
             "│   │   └── models.py\n"
             "│   ├── services/\n"
             "│   │   ├── ml_service.py ★ 584 lines\n"
             "│   │   ├── explain_service.py\n"
             "│   │   └── analysis_service.py\n"
             "│   ├── models/ (SQLAlchemy)\n"
             "│   ├── schemas/ (Pydantic)\n"
             "│   └── middleware/\n"
             "├── ml/\n"
             "│   ├── models/multitask_model.py\n"
             "│   ├── training/train.py\n"
             "│   └── data/ (13,792 samples)\n"
             "└── requirements.txt",
             font_size=10, color=LIGHT_GRAY, font_name="Consolas")

# Key Services
add_shape(slide, Inches(4.8), Inches(1.4), Inches(4), Inches(2.7), CARD_BG)
add_text_box(slide, Inches(5.1), Inches(1.55), Inches(3.4), Inches(0.4),
             "🧠 ML Service (Core)", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(5.1), Inches(2.05), Inches(3.4), Inches(1.8), [
    "584 lines of NLP logic",
    "28 emoji normalizations",
    "60+ sarcasm vocabulary words",
    "55+ misinformation vocabulary words",
    "18 compiled regex patterns",
    "Language detection algorithm",
], font_size=11, bullet_color=ACCENT_PURPLE)

add_shape(slide, Inches(9.1), Inches(1.4), Inches(3.8), Inches(2.7), CARD_BG)
add_text_box(slide, Inches(9.4), Inches(1.55), Inches(3.2), Inches(0.4),
             "🔍 Explain Service", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(9.4), Inches(2.05), Inches(3.2), Inches(1.8), [
    "Occlusion-based SHAP",
    "Token-level attribution",
    "Lexicon-based demo mode",
    "Supports both tasks",
    "Political ref detection",
], font_size=11, bullet_color=ACCENT_GREEN)

# Middleware
add_shape(slide, Inches(4.8), Inches(4.4), Inches(8.1), Inches(2.7), CARD_BG)
add_text_box(slide, Inches(5.1), Inches(4.55), Inches(7.5), Inches(0.4),
             "🛡️ Middleware & Infrastructure", font_size=16, color=ACCENT_ORANGE, bold=True)

mw_items = [
    ("CORS", "Allow cross-origin for frontend"),
    ("GZip", "Compress responses >1000 bytes"),
    ("Rate Limiter", "Sliding window: 10 req/60s per client"),
    ("Request Logger", "UUID request IDs + timing"),
    ("Cache (SHA-256)", "Deduplicate identical text analyses"),
    ("SQLite + Async", "Persistent storage via SQLAlchemy 2.0"),
]

for i, (title, desc) in enumerate(mw_items):
    col = i % 3
    row = i // 3
    x = Inches(5.1) + Inches(2.6) * col
    y = Inches(5.1) + Inches(0.8) * row
    add_text_box(slide, x, y, Inches(2.4), Inches(0.35),
                 f"● {title}", font_size=11, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(2.25), Inches(0.35),
                 desc, font_size=10, color=LIGHT_GRAY)

add_slide_number(slide, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Frontend Implementation
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Frontend Implementation", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(2.5), ACCENT_BLUE)

# Pages
pages = [
    ("🏠 Homepage", "Hero section with live demo\nAnimated stats (13,792 samples, F1)\n6 feature cards\nGradient animations", ACCENT_BLUE),
    ("🔬 Analyze Page", "Text input (500 char limit)\n4 example buttons\nResultCard with dual meters\nSHAP explain button\nJSON export", ACCENT_GREEN),
    ("📊 Dashboard", "StatsCards (total, rates, time)\nLine chart (14-day trend)\nRecent 5 analyses\nReal-time data via API", ACCENT_ORANGE),
    ("📋 History", "Paginated analysis table\nFilter: search, labels, dates\nRow expansion for details\nSorted by timestamp", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(pages):
    x = Inches(0.5) + Inches(3.15) * i
    add_shape(slide, x, Inches(1.4), Inches(2.9), Inches(3), CARD_BG)
    add_text_box(slide, x + Inches(0.2), Inches(1.55), Inches(2.5), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.05), Inches(2.5), Inches(2.1),
                 desc, font_size=11, color=LIGHT_GRAY)

# Key Components
add_shape(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(0.4),
             "Key UI Components", font_size=17, color=ACCENT_BLUE, bold=True)

components = [
    ("ResultCard", "Gradient accent bar, dual score\npanels, badges, descriptions", ACCENT_BLUE),
    ("ConfidenceMeter", "Animated SVG circular gauge\nGreen/Orange/Red auto-color", ACCENT_GREEN),
    ("ExplanationView", "3-tab: Sarcasm SHAP, Misinfo\nSHAP, Attention heatmap", ACCENT_ORANGE),
    ("AttentionHeatmap", "Canvas-rendered 20×20 matrix\nBlue color interpolation", ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(components):
    x = Inches(0.8) + Inches(3) * i
    add_text_box(slide, x, Inches(5.4), Inches(2.8), Inches(0.35),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, x, Inches(5.8), Inches(2.8), Inches(1),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Live Demo & Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Results — 14/14 Correct (100% Accuracy)", font_size=30, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3.5), ACCENT_GREEN)

results_data = [
    ["Input Text (Abbreviated)", "Sarcasm", "Misinfo", "Correct"],
    ["Modi ji ne toh desh ka bohot vikas 🙄", "98% ★", "1%", "✅"],
    ["Waah kya development, roads tut rahe", "98% ★", "1%", "✅"],
    ["50th meeting will definitely solve", "67%", "2%", "✅"],
    ["Hot water cures cancer! Share!", "1%", "98% ★", "✅"],
    ["Aaj ka weather bahut accha hai", "12%", "2%", "✅"],
    ["Highway inaugurated by CM", "2%", "2%", "✅"],
    ["10 crore spent, not a road fixed 😂🙄", "86% ★", "1%", "✅"],
    ["Cow urine COVID cure, forward to all", "1%", "98% ★", "✅"],
    ["Haan haan India is perfect", "92% ★", "1%", "✅"],
    ["Great speech PM Modi, inspiring", "12%", "2%", "✅"],
    ["5G se corona failta hai share karo", "2%", "98% ★", "✅"],
    ["Aaj match bahut exciting tha", "6%", "2%", "✅"],
    ["Budget allocated for education", "2%", "2%", "✅"],
    ["Scientists ne prove kiya 5G cancer", "1%", "98% ★", "✅"],
]

add_table(slide, Inches(0.3), Inches(1.3), len(results_data), 4,
          [Inches(6.5), Inches(1.7), Inches(1.7), Inches(1.2)], results_data)

add_slide_number(slide, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: Key Achievements & Future Work
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Key Achievements & Future Enhancements", font_size=30, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3.5), ACCENT_BLUE)

# Achievements
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.4),
             "✅ Key Achievements", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.5), [
    "100% accuracy on 14 diverse test cases",
    "Multi-task model: sarcasm + misinfo simultaneously",
    "13,792-sample dataset across 3 languages",
    "13 sarcasm rules + 8 misinformation rules",
    "28 emoji → semantic token normalizations",
    "SHAP explainability with token-level attribution",
    "Full-stack web app (FastAPI + Next.js 14)",
    "Real-time inference (<10ms per analysis)",
    "Dashboard with analytics & history tracking",
    "Docker-ready 5-service deployment",
], font_size=13, bullet_color=ACCENT_GREEN)

# Future
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.4),
             "🚀 Future Enhancements", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.5), [
    "Train XLM-RoBERTa on full dataset (GPU)",
    "Expand Tanglish (Tamil-English) patterns",
    "Real-time Twitter/X integration",
    "Browser extension for social media",
    "Multi-modal: meme/image analysis",
    "Regional languages: Marathi, Bengali, Kannada",
    "User feedback loop for model retraining",
    "Mobile responsive design",
    "Benchmark on SemEval & LIAR datasets",
    "Federated learning for privacy",
], font_size=13, bullet_color=ACCENT_ORANGE)

add_slide_number(slide, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: Thank You & Q&A
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

# Main text
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
             "Thank You!", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.3), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
             "Questions & Discussion", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Project summary
add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.6),
             "CodeMix NLP — Sarcasm & Misinformation Detection in Hinglish Text",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
             "XLM-RoBERTa  •  Multi-Task Learning  •  SHAP Explainability  •  FastAPI + Next.js",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Student info
add_shape(slide, Inches(4), Inches(5.8), Inches(5.3), Inches(1.0), CARD_BG)
add_text_box(slide, Inches(4), Inches(5.9), Inches(5.3), Inches(0.4),
             "Sumeet Sangwan", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4), Inches(6.3), Inches(5.3), Inches(0.4),
             "Vishwakarma University, Pune  |  B.Tech 3rd Year  |  NLP — Sem 2", 
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# Add slide transitions (fade) to ALL slides
# ═══════════════════════════════════════════════════════════════════════════════
for slide in prs.slides:
    transition = slide.element.makeelement(qn('p:transition'), {
        'spd': 'med',
        'advClick': '1'
    })
    fade = transition.makeelement(qn('p:fade'), {'thruBlk': '1'})
    transition.append(fade)
    slide.element.append(transition)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = os.path.join(
    r"c:\Users\Sumeet Sangwan\Desktop\Vishwakarma University\Vishwakarma University 3rd Year\Sem 2\NLP\Project",
    "CodeMix_NLP_Presentation.pptx"
)
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
